import io
import os
import random
import re
import shutil

import easyocr
import numpy as np
import requests  # 新增：用于调用智谱AI API
# 新增导入
import whisper
import win32com.client
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.picture import Picture

# 初始化 EasyOCR
reader = easyocr.Reader(['ch_sim', 'en'])  # 支持中英文

# 新增：初始化 Whisper 模型（使用 base 模型，速度快，准确率较高）
whisper_model = whisper.load_model("base")

def ocr_image_with_easyocr(img):
    """对PIL图片对象进行OCR识别，返回识别到的文本字符串"""
    ocr_result = reader.readtext(np.array(img), detail=0)
    texts = [item[1] for item in ocr_result if isinstance(item, (list, tuple)) and len(item) > 1]
    return '\n'.join(texts) if texts else ''

# 新增：音频转文字函数
def audio_to_text(audio_path):
    """将音频文件转为文字，支持mp3/wav/m4a等格式"""
    result = whisper_model.transcribe(audio_path)
    return result.get('text', '')

def ppt_to_text_list(ppt_path,change_pic=True):
    prs = Presentation(ppt_path)
    all_pages = []
    for slide in prs.slides:
        page_text = []
        for shape in slide.shapes:
            # 提取文本
            text = getattr(shape, "text", None)
            if text and text.strip():
                page_text.append(text)
            # 提取图片并OCR
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if change_pic and isinstance(shape, Picture):
                    image = shape.image
                    image_bytes = image.blob
                    image_stream = io.BytesIO(image_bytes)
                    img = Image.open(image_stream)
                    ocr_text = ocr_image_with_easyocr(img)
                    if ocr_text:
                        page_text.append(ocr_text)
        all_pages.append('\n'.join(page_text))
    return all_pages

# 新增：总结并生成问题的函数

def summarize_and_generate_questions(text_list, glm_api_key=None):
    """
    输入：字符串列表（如ppt_to_text_list的输出）
    输出：{
        'summary': 总结字符串,
        'questions': [
            {'question': str, 'options': [str, str, str, str]},
            ...
        ]
    }
    """
    if not text_list:
        return {'summary': '', 'questions': []}

    url = "https://open.bigmodel.cn/api/paas/v3/model-api/chatglm_std/invoke"
    headers = {
        "Authorization": f"Bearer {glm_api_key}",
        "Content-Type": "application/json"
    }
    prompt = (
        "你是一个智能助理。请根据以下内容：\n"
        + "\n".join(text_list) +
        "\n\n1. 只做简明扼要的总结（不超过150字）。\n"
        "2. 针对原始内容及相关常识，提出 5 个具有深度的关键问题，每个问题给出 4 个选项（只有一个正确，其余为干扰项）。\n"
        "3. 问题需基于内容推理判断得出，而非直接来源于原文表述，可适当扩展题目长度以保证完整性\n"
        "4. 问题题干不少于50字\n"
        "5. 问题需基于内容推理判断得出，而非直接来源于原文表述，可适当扩展题目长度以保证完整性。\n"
        "6. 不要出下面哪项是xxx之类直接得到答案的问题\n"
        "7. 每个问题请在选项后面输出“正确答案：A/B/C/D”和“解析：...”\n"
        "8. 基于"+str(random.random())+"随机出题\n"
        "输出格式要求：\n"
        "总结：...\n"
        "问题1：...\n"
        "A. ...\n"
        "B. ...\n"
        "C. ...\n"
        "D. ...\n"
        "正确答案：A\n"
        "解析：...\n"
        "问题2：...\n"
        "A. ...\n"
        "...\n"
        "正确答案：B\n"
        "解析：...\n"
    )
    data = {
        "model": "GLM-4-Plus",
        "prompt": prompt
    }
    response = requests.post(url, headers=headers, json=data)
    print("ChatGLM原始返回内容：", response.text)
    if response.status_code != 200:
        raise RuntimeError(f"ChatGLM API 调用失败: {response.text}")
    result_json = response.json()
    ai_output = result_json['data']['choices'][0]['content']

    # 先把字符串里的 \\n 替换成真正的换行符
    ai_output = ai_output.replace('\\n', '\n')
    lines = ai_output.strip().split('\n')
    summary = ''
    question_texts = []
    option_texts = []
    answer_indices = []
    explanations = []
    current_options = []
    for line in lines:
        line = line.strip()
        if not line:
            continue
        if line.startswith('总结：') or line.startswith('总结'):
            summary = line.replace('总结：', '').replace('总结', '').strip()
        elif re.match(r'^(\*+)?[#\s]*问题[1-9][0-9]*([：:：. ]|\*+)', line):
            if current_options:
                option_texts.extend(current_options)
                current_options = []
            q_text = re.sub(r'^(\*+)?[#\s]*问题[1-9][0-9]*([：:：. ]|\*+)+', '', line)
            question_texts.append(q_text)
        elif re.match(r'^[A-Da-d][.、:：\s]', line):
            option_text = re.sub(r'^[A-Da-d][.、:：\s]+', '', line)
            current_options.append(option_text)
        elif '正确答案' in line:
            match = re.search(r'[A-Da-d]', line)
            if match:
                idx = ord(match.group(0).upper()) - ord('A')
                answer_indices.append(str(idx))
        elif line.startswith('解析：') or line.startswith('解析:'):
            explanations.append(line.replace('解析：', '').replace('解析:', '').strip())
    if current_options:
        option_texts.extend(current_options)
    while len(answer_indices) < len(question_texts):
        answer_indices.append('0')
    while len(explanations) < len(question_texts):
        explanations.append('')
    questions_str = '\n'.join(question_texts)
    options_str = '\n'.join(option_texts)
    answer_indices_str = '\n'.join(answer_indices)
    explanations_str = '\n'.join(explanations)
    return {'summary': summary, 'questions_str': questions_str, 'options_str': options_str, 'answer_indices_str': answer_indices_str, 'explanations_str': explanations_str}

# 用法示例
def test_summarize_and_generate_questions(ppt_path, glm_api_key):
    """
    测试 summarize_and_generate_questions：
    1. 读取PPT内容转为文本列表
    2. 调用 summarize_and_generate_questions
    3. 打印所有问题和选项
    """
    text_list = ppt_to_text_list(ppt_path)
    result = summarize_and_generate_questions(text_list, glm_api_key=glm_api_key)
    print("总结：", result['summary'])
    print("问题列表：")
    print(result['questions_str'])
    print("选项列表：")
    print(result['options_str'])
    print("答案索引列表：")
    print(result['answer_indices_str'])
    print("解析列表：")
    print(result['explanations_str'])

def ppt_to_images(ppt_path, output_dir=None):
    """
    将PPT每一页渲染为图片，保存到当前目录或指定目录。
    图片命名为 page_1.png、page_2.png ...
    依赖：python-pptx、Pillow
    """
    import os
    from pptx import Presentation
    from PIL import Image, ImageDraw, ImageFont

    prs = Presentation(ppt_path)
    if output_dir is None:
        output_dir = os.getcwd()
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    # 尝试获取PPT页面尺寸
    width, height = prs.slide_width, prs.slide_height
    # 转为像素（1英寸=914400EMU, 1英寸=96像素）
    px_width = int(width * 96 / 914400)
    px_height = int(height * 96 / 914400)

    # 指定支持中文的字体路径（如simhei.ttf），如有需要请修改为你本机的字体文件
    font_path = "C:/Windows/Fonts/simhei.ttf"
    font = ImageFont.truetype(font_path, 24)

    for idx, slide in enumerate(prs.slides, 1):
        img = Image.new('RGB', (px_width, px_height), 'white')
        draw = ImageDraw.Draw(img)
        y = 20
        # 简单渲染文本内容
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text and text.strip():
                draw.text((20, y), text.strip(), fill='black', font=font)
                y += 40
        img_path = os.path.join(output_dir, f"page_{idx}.png")
        img.save(img_path)
    print(f"PPT每页图片已保存到: {output_dir}")

def ppt_to_images_with_office(ppt_path, output_dir=None):
    """
    使用PowerPoint自动化将PPT每页完整导出为图片，保留原始背景和格式。
    输出文件夹为PPT同名（不含扩展名），每页图片命名为1.JPG、2.JPG等。
    依赖：pip install pywin32，需本机安装Microsoft Office PowerPoint。
    """

    ppt_basename = os.path.splitext(os.path.basename(ppt_path))[0]
    parent_dir = os.path.dirname(ppt_path)
    # PowerPoint会自动在parent_dir下生成一个以文件名为名的目录
    ppt_basename=ppt_basename.split(".")[0]
    temp_output_dir = os.path.join(parent_dir, ppt_basename)
    if output_dir is None:
        output_dir = temp_output_dir
    # 先确保parent_dir存在
    if not os.path.exists(parent_dir):
        os.makedirs(parent_dir)
    ppt_app = win32com.client.Dispatch('PowerPoint.Application')
    ppt_app.Visible = 1
    presentation = ppt_app.Presentations.Open(ppt_path, WithWindow=False)
    # 17=ppSaveAsJPG，18=ppSaveAsPNG
    print(f"PowerPoint SaveAs目录: {temp_output_dir}")
    presentation.SaveAs(temp_output_dir, 17)
    presentation.Close()
    ppt_app.Quit()
    # 如果用户指定了output_dir且和PowerPoint生成的目录不同，则移动
    if os.path.abspath(output_dir) != os.path.abspath(temp_output_dir):
        if os.path.exists(output_dir):
            shutil.rmtree(output_dir)
        shutil.move(temp_output_dir, output_dir)
    # 重命名“幻灯片X.JPG”为“X.JPG”
    for fname in os.listdir(output_dir):
        match = re.match(r'幻灯片(\d+)\.JPG', fname, re.IGNORECASE)
        if match:
            new_name = f"{int(match.group(1))}.JPG"
            os.rename(os.path.join(output_dir, fname), os.path.join(output_dir, new_name))
    print(f'PPT每页图片已保存到: {output_dir}')
    return output_dir

if __name__ == "__main__":
    ppt_path = r"uploads\《Grandmaster+level+in+starcraft+II+using+multi-age.pptx"
    # result = ppt_to_text_list(ppt_path)
    # for i, page in enumerate(result, 1):
    #     print(f"--- Page {i} ---")
    #     print(page)
    #     print()
    # 新增：音频转文字演示
    audio_path = r"uploads\《Grandmaster+level+in+starcraft+II+using+multi-age.pptx"  # 替换为你的音频文件路径
    # audio_text = audio_to_text(audio_path)
    # print("--- Audio to Text ---")  
    # print(audio_text) 

    # 新增：测试 summarize_and_generate_questions
    glm_api_key = "a9205aba794f4f00acf33541eddbcd17.vqgIdbW54DlezvJh"  # TODO: 替换为你的智谱API Key
    test_summarize_and_generate_questions(ppt_path, glm_api_key)

    # 新增：测试 ppt_to_images_with_office
    # ppt_to_images_with_office(ppt_path) 